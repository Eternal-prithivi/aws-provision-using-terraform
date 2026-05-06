from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation("Truman_Phase2_Report.pptx")

BG = RGBColor(15, 17, 26)
CARD = RGBColor(22, 25, 38)
ACCENT = RGBColor(99, 102, 241)
GREEN = RGBColor(34, 197, 94)
AMBER = RGBColor(245, 158, 11)
WHITE = RGBColor(240, 240, 245)
MUTED = RGBColor(148, 163, 184)
BORDER = RGBColor(45, 50, 70)

def dark_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG

def add_box(slide, left, top, w, h, fill=CARD, border=BORDER):
    shp = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(w), Emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(1)
    return shp

def add_text(slide, left, top, w, h, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(w), Emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def accent_line(slide, left, top, w):
    shp = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(w), Emu(Inches(0.05)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = ACCENT
    shp.line.fill.background()

def footer(slide):
    add_text(slide, Inches(0.8), Inches(6.8), Inches(10), Inches(0.4), "Rooman Technologies  |  Internship Project — Phase II Build Report", 10, MUTED, False)

# ════════════════════════════════════════
# SLIDE 12 — CONCLUSION
# ════════════════════════════════════════
s12 = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(s12)
add_text(s12, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), "10 — Conclusion & Next Steps", 28, WHITE, True)
accent_line(s12, Inches(0.8), Inches(1.1), Inches(3))

# Summary box
add_box(s12, Inches(0.8), Inches(1.5), Inches(7.5), Inches(2.5))
add_text(s12, Inches(1.1), Inches(1.6), Inches(7), Inches(0.4), "What Was Built", 18, ACCENT, True)
built_items = [
    "✓  7 modular Terraform modules with feature-flag architecture",
    "✓  Dual-engine policy system (YAML + OPA) with 8 security rules",
    "✓  Full-stack Web Dashboard (FastAPI + Next.js) with RBAC",
    "✓  Automated drift detection with one-click remediation",
    "✓  CI/CD pipeline with 221 tests at 99% coverage",
    "✓  Slack-integrated approval workflows for team collaboration",
]
for i, item in enumerate(built_items):
    add_text(s12, Inches(1.1), Inches(2.1) + Emu(Inches(0.3) * i), Inches(7), Inches(0.3), item, 12, MUTED, False)

# Stats row
stats = [("7", "Terraform\nModules"), ("19", "API\nEndpoints"), ("221", "Tests\nPassing"), ("99%", "Code\nCoverage"), ("$0", "Monthly\nCost")]
for i, (num, label) in enumerate(stats):
    x = Inches(0.8) + Emu(Inches(2.4) * i)
    add_box(s12, x, Inches(4.3), Inches(2.2), Inches(1.1))
    add_text(s12, x, Inches(4.35), Inches(2.2), Inches(0.5), num, 28, ACCENT, True, PP_ALIGN.CENTER)
    add_text(s12, x, Inches(4.85), Inches(2.2), Inches(0.45), label, 10, MUTED, False, PP_ALIGN.CENTER)

# Next Steps
add_box(s12, Inches(8.8), Inches(1.5), Inches(4.0), Inches(2.5))
add_text(s12, Inches(9.1), Inches(1.6), Inches(3.5), Inches(0.4), "Next Steps (Phase 3)", 16, GREEN, True)
nexts = [
    "▸ Multi-region deployment support",
    "▸ Container orchestration (ECS/EKS)",
    "▸ Cost anomaly ML detection",
    "▸ SSO authentication integration",
    "▸ Production hardening & load testing",
]
for i, n in enumerate(nexts):
    add_text(s12, Inches(9.1), Inches(2.1) + Emu(Inches(0.35) * i), Inches(3.5), Inches(0.3), n, 11, MUTED, False)

# Thank you
add_text(s12, Inches(0.8), Inches(5.8), Inches(11.8), Inches(0.6),
    "The Smart AWS Infrastructure Provisioning System is fully operational, thoroughly tested,\nand ready for production integration.",
    14, WHITE, False, PP_ALIGN.CENTER)
add_text(s12, Inches(0.8), Inches(6.4), Inches(11.8), Inches(0.4), "Thank You", 20, ACCENT, True, PP_ALIGN.CENTER)
footer(s12)

prs.save("Truman_Phase2_Report.pptx")
print(f"ALL DONE! Total slides: {len(prs.slides)}")
